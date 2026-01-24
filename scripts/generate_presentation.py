#!/usr/bin/env python3
"""
Generate PowerPoint presentation for PharmaAnalytics project.
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create PharmaAnalytics PowerPoint presentation."""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    PRIMARY_BLUE = RGBColor(46, 134, 171)  # #2E86AB
    SECONDARY_GREEN = RGBColor(6, 167, 125)  # #06A77D
    ACCENT_ORANGE = RGBColor(241, 143, 1)  # #F18F01
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "PharmaAnalytics"
    subtitle.text = "Hospital Pharmacy Data Analytics Platform"
    
    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Project Overview"
    tf = content.text_frame
    tf.text = "What is PharmaAnalytics?"
    p = tf.add_paragraph()
    p.text = "• Comprehensive analytics platform for hospital pharmacy data"
    p = tf.add_paragraph()
    p.text = "• Process and analyze drug transaction data"
    p = tf.add_paragraph()
    p.text = "• Generate insights and visualizations"
    p = tf.add_paragraph()
    p.text = "• Forecast future drug demand"
    p = tf.add_paragraph()
    p.text = "• Perform time series diagnostics"
    
    # Slide 3: Architecture Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "System Architecture"
    tf = content.text_frame
    tf.text = "Technology Stack:"
    p = tf.add_paragraph()
    p.text = "• Backend: Flask (Python) REST API"
    p = tf.add_paragraph()
    p.text = "• Database: PostgreSQL 15"
    p = tf.add_paragraph()
    p.text = "• Task Queue: Celery with Redis"
    p = tf.add_paragraph()
    p.text = "• Containerization: Docker Compose"
    p = tf.add_paragraph()
    p.text = "• Data Processing: Polars & Pandas"
    p = tf.add_paragraph()
    p.text = "• Machine Learning: scikit-learn, XGBoost"
    
    # Slide 4: Technology Stack - Backend
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Backend Technology Stack"
    tf = content.text_frame
    tf.text = "Core Technologies:"
    p = tf.add_paragraph()
    p.text = "• Framework: Flask 3.0.0 (Python)"
    p = tf.add_paragraph()
    p.text = "• ORM: SQLAlchemy 2.0.23"
    p = tf.add_paragraph()
    p.text = "• Database: PostgreSQL 15"
    p = tf.add_paragraph()
    p.text = "• Migrations: Alembic 1.13.1"
    p = tf.add_paragraph()
    p.text = "• Validation: Pydantic 2.5.3"
    p = tf.add_paragraph()
    p.text = "• Background Tasks: Celery 5.3.4"
    p = tf.add_paragraph()
    p.text = "• Message Broker: Redis 7"
    
    # Slide 5: Technology Stack - Data & ML
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Data Processing & Machine Learning"
    tf = content.text_frame
    tf.text = "Data Processing:"
    p = tf.add_paragraph()
    p.text = "• Polars 0.20.0 (primary)"
    p = tf.add_paragraph()
    p.text = "• Pandas 2.1.4 (fallback)"
    p = tf.add_paragraph()
    p.text = "• PyArrow 14.0.1"
    p = tf.add_paragraph()
    p.text = "• Excel Support (openpyxl, xlrd)"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Machine Learning:"
    p = tf.add_paragraph()
    p.text = "• scikit-learn 1.4.0"
    p = tf.add_paragraph()
    p.text = "• XGBoost 2.0.3"
    p = tf.add_paragraph()
    p.text = "• NumPy 1.26.4"
    p = tf.add_paragraph()
    p.text = "• Joblib 1.3.2"
    
    # Slide 6: Core Modules Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "System Modules"
    tf = content.text_frame
    tf.text = "Four Main Modules:"
    p = tf.add_paragraph()
    p.text = "1. Data Ingestion Module - Excel file upload & processing"
    p = tf.add_paragraph()
    p.text = "2. Analytics Module - Business intelligence & reporting"
    p = tf.add_paragraph()
    p.text = "3. Diagnostics Module - Time series analysis"
    p = tf.add_paragraph()
    p.text = "4. Forecasting Module - ML-based demand prediction"
    
    # Slide 7: Module 1 - Data Ingestion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Data Ingestion Module"
    tf = content.text_frame
    tf.text = "Functionality:"
    p = tf.add_paragraph()
    p.text = "• Upload Excel files (.xlsx, .xls)"
    p = tf.add_paragraph()
    p.text = "• Validate and transform data"
    p = tf.add_paragraph()
    p.text = "• Track ingestion status"
    p = tf.add_paragraph()
    p.text = "• Error logging and review"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Features:"
    p = tf.add_paragraph()
    p.text = "• Asynchronous processing (Celery)"
    p = tf.add_paragraph()
    p.text = "• Status tracking (pending → processing → completed/failed)"
    p = tf.add_paragraph()
    p.text = "• Failed record storage for review"
    p = tf.add_paragraph()
    p.text = "• API Endpoint: /api/ingestion/*"
    
    # Slide 8: Module 2 - Analytics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Analytics Module"
    tf = content.text_frame
    tf.text = "Cost Analysis:"
    p = tf.add_paragraph()
    p.text = "• Sunburst charts, Trend analysis, Bubble charts"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Business Intelligence:"
    p = tf.add_paragraph()
    p.text = "• Top drugs by quantity"
    p = tf.add_paragraph()
    p.text = "• Drug demand time-series"
    p = tf.add_paragraph()
    p.text = "• Department performance"
    p = tf.add_paragraph()
    p.text = "• Year-over-year comparisons"
    p = tf.add_paragraph()
    p.text = "• Category analysis"
    p = tf.add_paragraph()
    p.text = "• Patient demographics"
    p = tf.add_paragraph()
    p.text = "• Summary statistics"
    p = tf.add_paragraph()
    p.text = "• API Endpoint: /api/analytics/*"
    
    # Slide 9: Module 3 - Diagnostics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Diagnostics Module"
    tf = content.text_frame
    tf.text = "Time Series Analysis:"
    p = tf.add_paragraph()
    p.text = "• Seasonality detection"
    p = tf.add_paragraph()
    p.text = "• Autocorrelation analysis"
    p = tf.add_paragraph()
    p.text = "• Outlier detection"
    p = tf.add_paragraph()
    p.text = "• Decomposition (trend, seasonal, residual)"
    p = tf.add_paragraph()
    p.text = "• Data profiling"
    p = tf.add_paragraph()
    p.text = "• Classification"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Performance:"
    p = tf.add_paragraph()
    p.text = "• Caching for faster results"
    p = tf.add_paragraph()
    p.text = "• API Endpoint: /api/diagnostics/*"
    
    # Slide 10: Module 4 - Forecasting
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Forecasting Module"
    tf = content.text_frame
    tf.text = "ML-Based Forecasting:"
    p = tf.add_paragraph()
    p.text = "• XGBoost algorithm"
    p = tf.add_paragraph()
    p.text = "• Domain-specific features"
    p = tf.add_paragraph()
    p.text = "• Model evaluation metrics"
    p = tf.add_paragraph()
    p.text = "• Forecast generation"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Use Cases:"
    p = tf.add_paragraph()
    p.text = "• Drug demand prediction"
    p = tf.add_paragraph()
    p.text = "• Inventory planning"
    p = tf.add_paragraph()
    p.text = "• Resource allocation"
    p = tf.add_paragraph()
    p.text = "• API Endpoint: /api/forecasting/*"
    
    # Slide 11: Data Model
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Database Schema"
    tf = content.text_frame
    tf.text = "Main Tables:"
    p = tf.add_paragraph()
    p.text = "• DrugTransaction - Transaction details, departments, patient info"
    p = tf.add_paragraph()
    p.text = "• DataIngestionLog - File upload tracking, status monitoring"
    p = tf.add_paragraph()
    p.text = "• DataIngestionError - Failed record storage, error details"
    
    # Slide 12: Data Schema Fields
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Data Input Schema"
    tf = content.text_frame
    tf.text = "Key Fields:"
    p = tf.add_paragraph()
    p.text = "• DOC (Document ID), DATE (Transaction date)"
    p = tf.add_paragraph()
    p.text = "• CODE (Drug code), ARTICLE (Drug name)"
    p = tf.add_paragraph()
    p.text = "• QTY (Quantity: negative = dispensed, positive = received)"
    p = tf.add_paragraph()
    p.text = "• U.P (Unit price), T.P (Total price)"
    p = tf.add_paragraph()
    p.text = "• C.R (Consuming department), C.S (Supplying department)"
    p = tf.add_paragraph()
    p.text = "• Patient info (Room, Bed, Date of birth)"
    
    # Slide 13: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Platform Features"
    tf = content.text_frame
    tf.text = "Key Capabilities:"
    p = tf.add_paragraph()
    p.text = "✓ Asynchronous Processing - Celery workers"
    p = tf.add_paragraph()
    p.text = "✓ Data Validation - Pydantic-based validation"
    p = tf.add_paragraph()
    p.text = "✓ Error Handling - Structured error tracking"
    p = tf.add_paragraph()
    p.text = "✓ Performance Optimization - Indexes, caching"
    p = tf.add_paragraph()
    p.text = "✓ RESTful API Design - Consistent endpoints"
    p = tf.add_paragraph()
    p.text = "✓ Dockerized Deployment - Containerized services"
    
    # Slide 14: Deployment Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Docker Compose Services"
    tf = content.text_frame
    tf.text = "Services:"
    p = tf.add_paragraph()
    p.text = "• postgres - PostgreSQL database (port 5433)"
    p = tf.add_paragraph()
    p.text = "• redis - Redis for Celery broker (port 6379)"
    p = tf.add_paragraph()
    p.text = "• backend - Flask API server (port 5000)"
    p = tf.add_paragraph()
    p.text = "• celery_worker - Background task processor"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Features:"
    p = tf.add_paragraph()
    p.text = "• Health checks for all services"
    p = tf.add_paragraph()
    p.text = "• Volume mounts for data persistence"
    p = tf.add_paragraph()
    p.text = "• Environment-based configuration"
    
    # Slide 15: Project Structure
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Codebase Organization"
    tf = content.text_frame
    tf.text = "Project Structure:"
    p = tf.add_paragraph()
    p.text = "• backend/ - Flask application (modules, database, shared)"
    p = tf.add_paragraph()
    p.text = "• docker/ - Dockerfiles"
    p = tf.add_paragraph()
    p.text = "• migrations/ - Database migrations"
    p = tf.add_paragraph()
    p.text = "• data/ - Data uploads and schema"
    p = tf.add_paragraph()
    p.text = "• scripts/ - Utility scripts"
    p = tf.add_paragraph()
    p.text = "• docs/ - Documentation"
    
    # Slide 16: Use Cases
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Who Uses PharmaAnalytics?"
    tf = content.text_frame
    tf.text = "Target Users:"
    p = tf.add_paragraph()
    p.text = "1. Hospital Administrators - Monitor usage and costs"
    p = tf.add_paragraph()
    p.text = "2. Pharmacy Managers - Track inventory and demand"
    p = tf.add_paragraph()
    p.text = "3. Financial Analysts - Cost analysis and trends"
    p = tf.add_paragraph()
    p.text = "4. Data Scientists - Forecasting and diagnostics"
    p = tf.add_paragraph()
    p.text = "5. Operations Teams - Department performance monitoring"
    
    # Slide 17: API Endpoints Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "REST API Overview"
    tf = content.text_frame
    tf.text = "API Endpoints:"
    p = tf.add_paragraph()
    p.text = "• /api/ingestion/* - File upload, status tracking"
    p = tf.add_paragraph()
    p.text = "• /api/analytics/* - Cost analysis, dashboards, reports"
    p = tf.add_paragraph()
    p.text = "• /api/diagnostics/* - Time series analysis"
    p = tf.add_paragraph()
    p.text = "• /api/forecasting/* - ML predictions"
    
    # Slide 18: Current Status
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Project Status"
    tf = content.text_frame
    tf.text = "Current State:"
    p = tf.add_paragraph()
    p.text = "✓ Backend API fully functional"
    p = tf.add_paragraph()
    p.text = "✓ Recent refactoring: Unified analytics module"
    p = tf.add_paragraph()
    p.text = "✓ Database migrations in place"
    p = tf.add_paragraph()
    p.text = "✓ ML forecasting capabilities implemented"
    p = tf.add_paragraph()
    p.text = "✓ Data ingestion pipeline operational"
    p = tf.add_paragraph()
    p.text = "✓ Comprehensive error handling"
    p = tf.add_paragraph()
    p.text = "✓ Docker deployment ready"
    
    # Slide 19: Benefits
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Platform Benefits"
    tf = content.text_frame
    tf.text = "Key Advantages:"
    p = tf.add_paragraph()
    p.text = "• Efficiency - Automated data processing"
    p = tf.add_paragraph()
    p.text = "• Insights - Real-time analytics"
    p = tf.add_paragraph()
    p.text = "• Forecasting - ML-powered predictions"
    p = tf.add_paragraph()
    p.text = "• Scalability - Docker-based deployment"
    p = tf.add_paragraph()
    p.text = "• Reliability - Robust error handling"
    p = tf.add_paragraph()
    p.text = "• Performance - Optimized queries and caching"
    
    # Slide 20: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You"
    subtitle.text = "Questions & Discussion"
    
    return prs


if __name__ == "__main__":
    print("Generating PowerPoint presentation...")
    prs = create_presentation()
    output_path = "PharmaAnalytics_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print("Note: Install python-pptx if not already installed: pip install python-pptx")

